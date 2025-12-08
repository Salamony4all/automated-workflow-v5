import os
from reportlab.lib.pagesizes import letter, A4
from reportlab.lib import colors
from reportlab.lib.units import inch
from reportlab.platypus import SimpleDocTemplate, Table, TableStyle, Paragraph, Spacer, Image as RLImage, PageBreak, Frame, PageTemplate
from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
from reportlab.lib.enums import TA_CENTER, TA_LEFT, TA_JUSTIFY
from reportlab.pdfgen import canvas
from datetime import datetime
import re
import logging
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from bs4 import BeautifulSoup

logger = logging.getLogger(__name__)

class PresentationGenerator:
    """Generate eye-catching technical presentations - 1 page per item"""
    
    def __init__(self):
        self.styles = getSampleStyleSheet()
        self.setup_custom_styles()
    
    def setup_custom_styles(self):
        """Setup custom paragraph styles for presentations"""
        self.title_style = ParagraphStyle(
            'PresentationTitle',
            parent=self.styles['Heading1'],
            fontSize=28,
            textColor=colors.HexColor('#1a365d'),
            spaceAfter=20,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        self.item_title_style = ParagraphStyle(
            'ItemTitle',
            parent=self.styles['Heading1'],
            fontSize=22,
            textColor=colors.HexColor('#1a365d'),
            spaceAfter=15,
            alignment=TA_CENTER,
            fontName='Helvetica-Bold'
        )
        
        self.spec_heading_style = ParagraphStyle(
            'SpecHeading',
            fontSize=14,
            textColor=colors.HexColor('#1a365d'),
            spaceAfter=8,
            fontName='Helvetica-Bold'
        )
        
        self.spec_text_style = ParagraphStyle(
            'SpecText',
            fontSize=11,
            textColor=colors.black,
            spaceAfter=6,
            alignment=TA_JUSTIFY
        )

    def _get_logo_path(self):
        candidates = [
            os.path.join('static', 'images', 'al-shaya-logo-white@2x.png'),  # White logo first for presentations
            os.path.join('static', 'images', 'AlShaya-Logo-color@2x.png'),
            os.path.join('static', 'images', 'LOGO.png')
        ]
        for p in candidates:
            if os.path.exists(p):
                return p
        return None

    def _get_white_logo_path(self):
        """Get white logo specifically for PPTX presentations"""
        white_logo = os.path.join('static', 'images', 'al-shaya-logo-white@2x.png')
        if os.path.exists(white_logo):
            return white_logo
        return self._get_logo_path()  # Fallback to default

    def _draw_header_footer(self, canv: canvas.Canvas, doc):
        """Draw properly placed header logo and footer website for presentation PDF."""
        page_width, page_height = doc.pagesize
        gold = colors.HexColor('#d4af37')
        dark = colors.HexColor('#1a365d')
        
        # Header gold line
        canv.setStrokeColor(gold)
        canv.setLineWidth(2)
        canv.line(doc.leftMargin, page_height - 40, page_width - doc.rightMargin, page_height - 40)
        
        # Logo centered in header - larger and more visible
        logo = self._get_logo_path()
        if logo and os.path.exists(logo):
            try:
                w, h = 130, 46  # Larger logo
                x = (page_width - w) / 2  # Center horizontally
                y = page_height - 38
                canv.drawImage(logo, x, y, width=w, height=h, preserveAspectRatio=True, mask='auto')
            except Exception:
                pass
        
        # Footer with gold line and website centered
        canv.setStrokeColor(gold)
        canv.setLineWidth(2)
        canv.line(doc.leftMargin, doc.bottomMargin + 15, page_width - doc.rightMargin, doc.bottomMargin + 15)
        
        canv.setFillColor(dark)
        canv.setFont('Helvetica', 10)
        footer_text = 'https://alshayaenterprises.com'
        canv.drawCentredString(page_width / 2, doc.bottomMargin + 5, footer_text)
    
    def generate(self, file_id, session, format_type='pdf'):
        """
        Generate technical presentation with 1 page/slide per item
        Always generates PPTX first, then converts to PDF if needed
        Args:
            file_id: The file ID
            session: Flask session
            format_type: 'pdf' or 'pptx'
        Returns: path to generated file (PDF or PPTX), and stores PPTX path in session
        """
        # Get file info and extracted data
        uploaded_files = session.get('uploaded_files', [])
        file_info = None
        
        for f in uploaded_files:
            if f['id'] == file_id:
                file_info = f
                break
        
        if not file_info:
            raise Exception('File not found. Please upload and extract a file first.')
        
        # Check if this is multi-budget and get product selections
        is_multibudget = file_info.get('multibudget', False)
        product_selections = file_info.get('product_selections', []) if is_multibudget else []
        tier = file_info.get('tier', 'budgetary') if is_multibudget else None
        
        # Get costed data (preferred) or stitched table or extraction result
        if 'costed_data' in file_info:
            items = self.parse_items_from_costed_data(file_info['costed_data'], session, file_id, 
                                                      is_multibudget, product_selections, tier)
        elif 'stitched_table' in file_info:
            items = self.parse_items_from_stitched_table(file_info['stitched_table'], session, file_id,
                                                         is_multibudget, product_selections, tier)
        elif 'extraction_result' in file_info:
            items = self.parse_items_from_extraction(file_info['extraction_result'], session, file_id)
        else:
            raise Exception('No data available. Please extract tables first.')
        
        if not items:
            raise Exception('No items found in the table.')
        
        # Create output directory
        session_id = session['session_id']
        output_dir = os.path.join('outputs', session_id, 'presentations')
        os.makedirs(output_dir, exist_ok=True)
        
        # Generate file based on format
        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        
        # ALWAYS generate PPTX first (will be available for download)
        pptx_file = os.path.join(output_dir, f'presentation_{file_id}_{timestamp}.pptx')
        self.generate_pptx(items, pptx_file)
        
        # Store PPTX path in file_info for later download
        file_info['presentation_pptx'] = pptx_file
        
        if format_type == 'pptx':
            output_file = pptx_file
        else:  # pdf - convert PPTX to PDF to maintain same layout
            output_file = os.path.join(output_dir, f'presentation_{file_id}_{timestamp}.pdf')
            self.convert_pptx_to_pdf(pptx_file, output_file)
        
        return output_file
    
    def parse_items_from_costed_data(self, costed_data, session, file_id, is_multibudget=False, 
                                     product_selections=None, tier=None):
        """Parse items from costed table data"""
        items = []
        session_id = session.get('session_id', '')
        product_selections = product_selections or []
        
        for table in costed_data.get('tables', []):
            headers = [h for h in table.get('headers', []) if str(h).lower() not in ['action', 'actions', 'product selection', 'productselection']]
            
            for row_idx, row in enumerate(table.get('rows', [])):
                # Log all available columns for debugging
                logger.info(f"Row headers: {list(row.keys())}")
                
                # Find description column - for multi-budget, prioritize Brand Description
                description = ''
                raw_description = ''
                description_found = False
                
                # For multi-budget: Priority 1 - Brand Description from costed table
                if is_multibudget:
                    for h in headers:
                        h_str = str(h).lower() if h else ''
                        if 'brand description' in h_str or (h_str == 'brand description'):
                            raw_description = row.get(h, '')
                            description = self.strip_html(raw_description)
                            if description and description.strip() and 'no description' not in description.lower():
                                logger.info(f"Found BRAND DESCRIPTION column '{h}' (length: {len(description)}): {description[:150]}...")
                                description_found = True
                                break
                
                # Priority order: Brand Description (multi-budget) > DESCRIPTION > Item > Product
                if not description_found:
                    for h in headers:
                        h_str = str(h).lower() if h else ''
                        # Check for description column (most detailed) - but skip Brand Description if already checked
                        if ('descript' in h_str or 'discript' in h_str) and 'brand' not in h_str:
                            raw_description = row.get(h, '')
                            description = self.strip_html(raw_description)
                            logger.info(f"Found DESCRIPTION column '{h}' (length: {len(description)}): {description[:150]}...")
                            description_found = True
                            break
                
                # If no description column found, try item or product columns
                if not description_found:
                    for h in headers:
                        h_str = str(h).lower() if h else ''
                        if 'item' in h_str or 'product' in h_str:
                            raw_description = row.get(h, '')
                            description = self.strip_html(raw_description)
                            logger.info(f"Found ITEM/PRODUCT column '{h}' (length: {len(description)}): {description[:150]}...")
                            break
                
                # Find quantity
                qty = ''
                unit = ''
                for h in headers:
                    h_str = str(h).lower() if h else ''
                    if 'qty' in h_str or 'quantity' in h_str:
                        qty = self.strip_html(row.get(h, ''))
                    if 'unit' in h_str and 'rate' not in h_str:
                        unit = self.strip_html(row.get(h, ''))
                
                # Find pricing
                unit_rate = ''
                total = ''
                for h in headers:
                    h_str = str(h).lower() if h else ''
                    if 'rate' in h_str or 'price' in h_str:
                        unit_rate = self.strip_html(row.get(h, ''))
                    if 'total' in h_str or 'amount' in h_str:
                        total = self.strip_html(row.get(h, ''))
                
                # Find reference image(s) from table - for multi-budget, this will be small reference image
                reference_image_paths = []
                image_paths = []  # For non-multi-budget, regular images
                selected_product_image = None  # For multi-budget: Brand Image from costed table
                
                for h in headers:
                    h_str = str(h).lower() if h else ''
                    cell_value = row.get(h, '')
                    
                    if is_multibudget:
                        # For multi-budget: Priority 1 - Brand Image from costed table
                        if 'brand image' in h_str or h_str == 'brand image':
                            if self.contains_image(cell_value):
                                paths = self.extract_all_image_paths(cell_value, session_id, file_id)
                                if paths:
                                    selected_product_image = paths[0]  # Use first Brand Image
                                    logger.info(f"Found BRAND IMAGE column '{h}' for multi-budget")
                        
                        # For multi-budget: look for indicative/reference image (not Brand Image)
                        elif ('indicative' in h_str and 'image' in h_str) or ('image' in h_str and 'brand' not in h_str and 'product' not in h_str):
                            if self.contains_image(cell_value):
                                paths = self.extract_all_image_paths(cell_value, session_id, file_id)
                                if paths:
                                    reference_image_paths.extend(paths)
                    else:
                        # For non-multi-budget: look for any image column
                        if 'image' in h_str:
                            if self.contains_image(cell_value):
                                paths = self.extract_all_image_paths(cell_value, session_id, file_id)
                                if paths:
                                    image_paths.extend(paths)
                
                # For multi-budget: download Brand Image if it's a URL
                if is_multibudget and selected_product_image and selected_product_image.startswith('http'):
                    from utils.image_helper import download_image
                    cached_path = download_image(selected_product_image)
                    if cached_path:
                        selected_product_image = cached_path
                
                # Use Brand Description and Brand Image from costed table for multi-budget
                final_description = description  # Already extracted Brand Description above for multi-budget
                final_image_paths = [selected_product_image] if (is_multibudget and selected_product_image) else (image_paths if image_paths else [])
                
                if final_description:  # Only add if we have a description
                    item = {
                        'description': final_description,
                        'qty': qty,
                        'unit': unit,
                        'unit_rate': unit_rate,
                        'total': total,
                        'image_path': final_image_paths[0] if final_image_paths else None,  # Selected product image (big)
                        'image_paths': final_image_paths,  # Selected product images
                        'reference_image_path': reference_image_paths[0] if reference_image_paths else None,  # Reference image (small) for multi-budget
                        'reference_image_paths': reference_image_paths,  # All reference images
                        'is_multibudget': is_multibudget,  # Flag to indicate multi-budget
                        'brand': self.extract_brand(final_description),
                        'specifications': self.extract_specifications(final_description)
                    }
                    items.append(item)
        
        return items
    
    def parse_items_from_stitched_table(self, stitched_table, session, file_id, is_multibudget=False,
                                        product_selections=None, tier=None):
        """Parse items from stitched HTML table data"""
        items = []
        session_id = session.get('session_id', '')
        product_selections = product_selections or []
        
        # Parse the HTML
        html_content = stitched_table.get('html', '')
        soup = BeautifulSoup(html_content, 'html.parser')
        
        # Find the table
        table = soup.find('table')
        if not table:
            logger.error("No table found in stitched HTML")
            return items
        
        # Get headers
        headers = []
        header_row = table.find('tr')
        if header_row:
            for th in header_row.find_all(['th', 'td']):
                header_text = th.get_text(strip=True).lower()
                # Exclude Product Selection and Actions columns
                if header_text not in ['action', 'actions', 'product selection', 'productselection']:
                    headers.append(header_text)
        
        logger.info(f"Found headers: {headers}")
        logger.info(f"Checking for description header in: {headers}")
        
        # Get data rows (skip header row)
        rows = table.find_all('tr')[1:]  # Skip first row (headers)
        logger.info(f"Found {len(rows)} data rows")
        
        for row_idx, row in enumerate(rows):
            cells = row.find_all('td')
            
            # Build row dict, skipping Product Selection and Actions cells
            row_data = {}
            col_idx = 0
            for i, cell in enumerate(cells):
                # Skip Product Selection and Actions cells
                if cell.find(class_='product-selection-dropdowns') or cell.find('button'):
                    continue
                text = cell.get_text(strip=True).lower()
                if 'product selection' in text or 'actions' in text:
                    continue
                
                if col_idx < len(headers):
                    # Check if cell contains image
                    img = cell.find('img')
                    if img:
                        row_data[headers[col_idx]] = str(cell)  # Keep HTML with image
                    else:
                        row_data[headers[col_idx]] = cell.get_text(strip=True)
                    col_idx += 1
            
            # Extract fields
            description = ''
            logger.debug(f"Looking for description in headers: {headers}")
            logger.debug(f"Row data keys: {list(row_data.keys())}")
            
            description_found = False
            # Priority: DESCRIPTION column first
            for h in headers:
                h_str = str(h).lower() if h else ''
                if 'descript' in h_str or 'discript' in h_str:
                    description = self.strip_html(row_data.get(h, ''))
                    logger.info(f"Stitched: Found DESCRIPTION in '{h}' (length: {len(description)}): {description[:150]}...")
                    description_found = True
                    break
            
            # If no description found, try item or product columns
            if not description_found:
                for h in headers:
                    h_str = str(h).lower() if h else ''
                    if 'item' in h_str or 'product' in h_str:
                        description = self.strip_html(row_data.get(h, ''))
                        logger.info(f"Stitched: Found ITEM/PRODUCT in '{h}' (length: {len(description)}): {description[:150]}...")
                        break
            
            if not description:
                logger.warning(f"No description found in row. Headers: {headers}, Row data keys: {list(row_data.keys())}")
                continue
            
            # Find quantity
            qty = ''
            unit = ''
            for h in headers:
                h_str = str(h).lower() if h else ''
                if 'qty' in h_str or 'quantity' in h_str:
                    qty = self.strip_html(row_data.get(h, ''))
                if 'unit' in h_str and 'rate' not in h_str and 'price' not in h_str:
                    unit = self.strip_html(row_data.get(h, ''))
            
            # Find pricing
            unit_rate = ''
            total = ''
            for h in headers:
                h_str = str(h).lower() if h else ''
                if ('rate' in h_str or 'price' in h_str) and 'unit' in h_str:
                    unit_rate = self.strip_html(row_data.get(h, ''))
                if 'total' in h_str or 'amount' in h_str:
                    total = self.strip_html(row_data.get(h, ''))
            
            # Find reference image(s) from table - for multi-budget, this will be small reference image
            reference_image_paths = []
            image_paths = []  # For non-multi-budget, regular images
            selected_product_image = None  # For multi-budget: Brand Image from costed table
            
            for h in headers:
                h_str = str(h).lower() if h else ''
                cell_value = row_data.get(h, '')
                
                if is_multibudget:
                    # For multi-budget: Priority 1 - Brand Image from costed table
                    if 'brand image' in h_str or h_str == 'brand image':
                        if self.contains_image(str(cell_value)):
                            paths = self.extract_all_image_paths(str(cell_value), session_id, file_id)
                            if paths:
                                selected_product_image = paths[0]  # Use first Brand Image
                                logger.info(f"Stitched: Found BRAND IMAGE column '{h}' for multi-budget")
                    
                    # For multi-budget: look for indicative/reference image (not Brand Image)
                    elif ('indicative' in h_str and 'image' in h_str) or ('image' in h_str and 'brand' not in h_str and 'product' not in h_str):
                        if self.contains_image(str(cell_value)):
                            paths = self.extract_all_image_paths(str(cell_value), session_id, file_id)
                            if paths:
                                reference_image_paths.extend(paths)
                else:
                    # For non-multi-budget: look for any image column
                    if 'image' in h_str:
                        if self.contains_image(str(cell_value)):
                            paths = self.extract_all_image_paths(str(cell_value), session_id, file_id)
                            if paths:
                                image_paths.extend(paths)
            
            # For multi-budget: download Brand Image if it's a URL
            if is_multibudget and selected_product_image and selected_product_image.startswith('http'):
                from utils.image_helper import download_image
                cached_path = download_image(selected_product_image)
                if cached_path:
                    selected_product_image = cached_path
            
            # Use Brand Description and Brand Image from costed table for multi-budget
            final_description = description  # Already extracted Brand Description above for multi-budget
            final_image_paths = [selected_product_image] if (is_multibudget and selected_product_image) else (reference_image_paths if reference_image_paths else [])
            
            item = {
                'description': final_description,
                'qty': qty,
                'unit': unit,
                'unit_rate': unit_rate,
                'total': total,
                'image_path': final_image_paths[0] if final_image_paths else None,  # Selected product image (big)
                'image_paths': final_image_paths,  # Selected product images
                'reference_image_path': reference_image_paths[0] if reference_image_paths else None,  # Reference image (small) for multi-budget
                'reference_image_paths': reference_image_paths,  # All reference images
                'is_multibudget': is_multibudget,  # Flag to indicate multi-budget
                'brand': self.extract_brand(final_description),
                'specifications': self.extract_specifications(final_description)
            }
            items.append(item)
        
        logger.info(f"Parsed {len(items)} items from stitched table")
        return items
    
    def strip_html(self, text):
        """Strip HTML tags from text but preserve all text content including from image alt tags"""
        text = str(text)
        
        # Extract text from img alt attributes before removing tags
        img_alts = re.findall(r'<img[^>]*alt=["\']([^"\']+)["\']', text, re.IGNORECASE)
        
        # Remove HTML tags
        text = re.sub(r'<[^>]+>', ' ', text)
        
        # Decode HTML entities
        text = text.replace('&nbsp;', ' ').replace('&amp;', '&').replace('&lt;', '<').replace('&gt;', '>')
        
        # Clean up multiple spaces and normalize whitespace
        text = re.sub(r'\s+', ' ', text).strip()
        
        return text
    
    def contains_image(self, cell_value):
        """Check if cell contains an image reference"""
        return '<img' in str(cell_value).lower() or 'img_in_' in str(cell_value).lower()
    
    def extract_all_image_paths(self, cell_value, session_id, file_id):
        """Extract ALL image paths from cell value (supports multiple images)"""
        image_paths = []
        try:
            # Find all src="..." patterns
            matches = re.findall(r'src=["\']([^"\']+)["\']', str(cell_value))
            
            for img_path in matches:
                img_path = img_path.lstrip('/')
                
                # Handle URLs (http/https)
                if img_path.startswith('http://') or img_path.startswith('https://'):
                    image_paths.append(img_path)
                    continue
                
                # Handle local paths - if already starts with outputs, return as-is
                if img_path.startswith('outputs'):
                    image_paths.append(img_path)
                    continue
                
                # Ensure all parts are strings
                if isinstance(session_id, (list, tuple)):
                    session_id = session_id[0] if session_id else ''
                if isinstance(file_id, (list, tuple)):
                    file_id = file_id[0] if file_id else ''
                if isinstance(img_path, (list, tuple)):
                    img_path = img_path[0] if img_path else ''
                
                # Check if it's a relative path that needs to be joined
                full_path = os.path.join('outputs', str(session_id), str(file_id), str(img_path))
                if os.path.exists(full_path):
                    image_paths.append(full_path)
                else:
                    # Also try without the session_id/file_id prefix
                    if os.path.exists(str(img_path)):
                        image_paths.append(str(img_path))
                    else:
                        image_paths.append(full_path)  # Return even if doesn't exist yet
            
            # Also look for any image path pattern in imgs/ folder (more flexible regex)
            if 'imgs/' in str(cell_value) and not matches:
                img_matches = re.findall(r'(imgs/[^"\s<>]+\.(jpg|png|jpeg|gif|webp))', str(cell_value), re.IGNORECASE)
                for img_relative_path, _ in img_matches:
                    if isinstance(session_id, (list, tuple)):
                        session_id = session_id[0] if session_id else ''
                    if isinstance(file_id, (list, tuple)):
                        file_id = file_id[0] if file_id else ''
                    full_path = os.path.join('outputs', str(session_id), str(file_id), str(img_relative_path))
                    image_paths.append(full_path)
                    
        except Exception as e:
            logger.error(f"Error extracting image paths: {e}")
        
        return image_paths if image_paths else None
    
    def extract_image_path(self, cell_value, session_id, file_id):
        """Extract first image path from cell value (for backward compatibility)"""
        paths = self.extract_all_image_paths(cell_value, session_id, file_id)
        return paths[0] if paths else None
    
    def generate_pdf(self, items, output_file):
        """Generate PDF presentation"""
        doc = SimpleDocTemplate(output_file, pagesize=A4, 
                                topMargin=1.0*inch, bottomMargin=0.8*inch,
                                leftMargin=0.75*inch, rightMargin=0.75*inch)
        story = []
        
        # Cover page
        story.extend(self.create_cover_page())
        story.append(PageBreak())
        
        # Create one page per item
        for idx, item in enumerate(items):
            story.extend(self.create_item_page_pdf(item, idx + 1))
            if idx < len(items) - 1:
                story.append(PageBreak())
        
        # Build PDF with header/footer
        doc.build(story, onFirstPage=self._draw_header_footer, onLaterPages=self._draw_header_footer)
    
    def generate_pptx(self, items, output_file):
        """Generate PowerPoint presentation"""
        prs = Presentation()
        prs.slide_width = Inches(10)
        prs.slide_height = Inches(7.5)
        
        # Add title slide
        self.create_title_slide_pptx(prs)
        
        # Add one slide per item
        for idx, item in enumerate(items):
            self.create_item_slide_pptx(prs, item, idx + 1)
        
        prs.save(output_file)
    
    def convert_pptx_to_pdf(self, pptx_file, pdf_file):
        """Convert PPTX to PDF using LibreOffice (cross-platform) or PowerPoint (Windows only)"""
        import platform
        import subprocess
        import time
        
        pptx_abs = os.path.abspath(pptx_file)
        pdf_abs = os.path.abspath(pdf_file)
        output_dir = os.path.dirname(pdf_abs)
        
        logger.info(f"Converting PPTX to PDF: {pptx_abs} -> {pdf_abs}")
        
        try:
            # Try LibreOffice first (works on Linux/Railway)
            libreoffice_cmds = [
                'libreoffice',
                'soffice',
                '/usr/bin/libreoffice',
                '/usr/bin/soffice'
            ]
            
            for cmd in libreoffice_cmds:
                try:
                    # LibreOffice headless conversion
                    result = subprocess.run(
                        [cmd, '--headless', '--convert-to', 'pdf', '--outdir', output_dir, pptx_abs],
                        capture_output=True,
                        text=True,
                        timeout=60
                    )
                    
                    if result.returncode == 0:
                        # LibreOffice creates PDF with same name as PPTX
                        expected_pdf = os.path.join(output_dir, os.path.splitext(os.path.basename(pptx_abs))[0] + '.pdf')
                        
                        # Wait for file to be created
                        for _ in range(20):
                            if os.path.exists(expected_pdf) and os.path.getsize(expected_pdf) > 1000:
                                # Rename if necessary
                                if expected_pdf != pdf_abs:
                                    if os.path.exists(pdf_abs):
                                        os.remove(pdf_abs)
                                    os.rename(expected_pdf, pdf_abs)
                                
                                logger.info(f"PDF created successfully with LibreOffice: {pdf_abs}")
                                return
                            time.sleep(0.5)
                        
                        logger.warning(f"LibreOffice conversion completed but PDF not found: {expected_pdf}")
                    else:
                        logger.debug(f"LibreOffice command failed: {result.stderr}")
                        
                except FileNotFoundError:
                    continue  # Try next command
                except Exception as e:
                    logger.debug(f"LibreOffice attempt with {cmd} failed: {e}")
                    continue
            
            # Fallback to PowerPoint on Windows
            if platform.system() == 'Windows':
                try:
                    import comtypes.client
                    
                    logger.info("Attempting PowerPoint COM conversion on Windows")
                    
                    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                    powerpoint.Visible = 1
                    
                    presentation = powerpoint.Presentations.Open(pptx_abs)
                    presentation.SaveAs(pdf_abs, 32)  # 32 = ppSaveAsPDF
                    presentation.Close()
                    powerpoint.Quit()
                    
                    # Wait for file
                    for _ in range(20):
                        if os.path.exists(pdf_abs) and os.path.getsize(pdf_abs) > 1000:
                            logger.info(f"PDF created successfully with PowerPoint: {pdf_abs}")
                            return
                        time.sleep(0.5)
                    
                except Exception as e:
                    logger.error(f"PowerPoint COM conversion failed: {e}")
            
            # If all methods fail, raise error
            raise Exception("Could not convert PPTX to PDF. LibreOffice/PowerPoint not available or conversion failed.")
            
        except Exception as e:
            logger.error(f"Error converting PPTX to PDF: {e}")
            import traceback
            logger.error(traceback.format_exc())
            raise Exception(f"Could not convert presentation to PDF. Error: {str(e)}")
    
    def create_title_slide_pptx(self, prs):
        """Create PowerPoint title slide with enhanced design"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Background - Navy blue header box (increased height to contain full logo)
        header_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0), 
            Inches(10), Inches(3.2)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RGBColor(26, 54, 93)  # Navy blue
        header_shape.line.fill.background()
        
        # Gold accent bar (moved down to be below the logo)
        accent_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(3.2),
            Inches(10), Inches(0.15)
        )
        accent_shape.fill.solid()
        accent_shape.fill.fore_color.rgb = RGBColor(212, 175, 55)  # Gold
        accent_shape.line.fill.background()
        
        # Logo centered in header with more vertical space - use white logo
        logo = self._get_white_logo_path()
        if logo and os.path.exists(logo):
            try:
                slide.shapes.add_picture(logo, Inches(3.5), Inches(0.8), width=Inches(3))
            except Exception:
                pass
        
        # Title with navy background (no company name - removed as requested)
        title_box = slide.shapes.add_textbox(Inches(1), Inches(3.9), Inches(8), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "TECHNICAL PROPOSAL"
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(48)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(26, 54, 93)  # Navy
        title_p.alignment = PP_ALIGN.CENTER
        
        # Subtitle
        subtitle_box = slide.shapes.add_textbox(Inches(2), Inches(5.0), Inches(6), Inches(0.6))
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = "Furniture, Fixtures & Equipment"
        subtitle_p = subtitle_frame.paragraphs[0]
        subtitle_p.font.size = Pt(26)
        subtitle_p.font.color.rgb = RGBColor(100, 116, 139)  # Gray
        subtitle_p.alignment = PP_ALIGN.CENTER
        
        # Date
        date_box = slide.shapes.add_textbox(Inches(3), Inches(6.2), Inches(4), Inches(0.5))
        date_frame = date_box.text_frame
        date_frame.text = datetime.now().strftime('%B %d, %Y')
        date_p = date_frame.paragraphs[0]
        date_p.font.size = Pt(18)
        date_p.font.color.rgb = RGBColor(71, 85, 105)  # Dark gray
        date_p.alignment = PP_ALIGN.CENTER
        
        # Footer with website
        footer_box = slide.shapes.add_textbox(Inches(2.5), Inches(6.8), Inches(5), Inches(0.4))
        footer_frame = footer_box.text_frame
        footer_frame.text = "https://alshayaenterprises.com"
        footer_p = footer_frame.paragraphs[0]
        footer_p.font.size = Pt(14)
        footer_p.font.color.rgb = RGBColor(212, 175, 55)  # Gold
        footer_p.alignment = PP_ALIGN.CENTER
    
    def create_item_slide_pptx(self, prs, item, page_num):
        """Create PowerPoint slide for one item with enhanced design"""
        slide_layout = prs.slide_layouts[6]  # Blank layout
        slide = prs.slides.add_slide(slide_layout)
        
        # Header bar with navy blue background - increased height
        header_shape = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(0),
            Inches(10), Inches(1.1)
        )
        header_shape.fill.solid()
        header_shape.fill.fore_color.rgb = RGBColor(26, 54, 93)  # Navy blue
        header_shape.line.fill.background()
        
        # Gold accent line under header
        accent_line = slide.shapes.add_shape(
            1,  # Rectangle
            Inches(0), Inches(1.1),
            Inches(10), Inches(0.08)
        )
        accent_line.fill.solid()
        accent_line.fill.fore_color.rgb = RGBColor(212, 175, 55)  # Gold
        accent_line.line.fill.background()
        
        # Small white logo top-right in header
        logo = self._get_white_logo_path()
        if logo and os.path.exists(logo):
            try:
                slide.shapes.add_picture(logo, Inches(8.2), Inches(0.2), width=Inches(1.5))
            except Exception:
                pass
        
        # For multi-budget: Add small reference image in upper left below gold line
        is_multibudget = item.get('is_multibudget', False)
        reference_image_path = item.get('reference_image_path')
        if is_multibudget and reference_image_path:
            try:
                # Download if URL
                if reference_image_path.startswith('http'):
                    from utils.image_helper import download_image
                    cached_path = download_image(reference_image_path)
                    if cached_path:
                        reference_image_path = cached_path
                
                if reference_image_path and os.path.exists(reference_image_path):
                    # Small reference image in upper left (below gold line with spacing)
                    ref_img = slide.shapes.add_picture(reference_image_path, Inches(0.5), Inches(1.35), width=Inches(1.2), height=Inches(0.9))
                    
                    # Add "Reference Image" label below the image
                    ref_label_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.25), Inches(1.2), Inches(0.2))
                    ref_label_frame = ref_label_box.text_frame
                    ref_label_frame.text = "Reference Image"
                    ref_label_p = ref_label_frame.paragraphs[0]
                    ref_label_p.font.size = Pt(8)
                    ref_label_p.font.color.rgb = RGBColor(100, 100, 100)  # Gray
                    ref_label_p.alignment = PP_ALIGN.CENTER
            except Exception as e:
                logger.warning(f"Could not add reference image: {e}")
        
        # Title in header - show item number and short title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(7.2), Inches(0.6))
        title_frame = title_box.text_frame
        title_frame.word_wrap = True
        # Extract first line or first 60 chars for title
        first_line = item['description'].split('\n')[0] if '\n' in item['description'] else item['description']
        title_text = f"Item {page_num}: {first_line[:60]}" + ("..." if len(first_line) > 60 else "")
        title_frame.text = title_text
        title_p = title_frame.paragraphs[0]
        title_p.font.size = Pt(18 if len(first_line) > 50 else 20)
        title_p.font.bold = True
        title_p.font.color.rgb = RGBColor(255, 255, 255)  # White text on navy
        
        # Images (left side) - support multiple images, adjusted position to account for taller header
        # For multi-budget: use selected product images (bigger), not reference images
        image_paths = item.get('image_paths', [item.get('image_path')] if item.get('image_path') else [])
        # Filter out None values
        image_paths = [p for p in image_paths if p]
        
        # Adjust position if multi-budget (to account for reference image in upper left)
        image_area_y = Inches(1.8) if not is_multibudget else Inches(2.5)  # Lower if reference image present
        
        logger.info(f"Item {page_num}: Found {len(image_paths)} image(s) to display (multibudget: {is_multibudget})")
        
        if image_paths:
            try:
                from PIL import Image as PILImage
                
                # Define image area - adjust for multi-budget
                area_x = Inches(0.5)
                area_y = image_area_y  # Use adjusted Y position
                area_w = Inches(4.5)
                area_h = Inches(4.0) if is_multibudget else Inches(4.5)  # Slightly smaller if reference image present
                
                num_images = min(len(image_paths), 6) # Max 6 images
                image_paths = image_paths[:num_images]
                
                # Determine grid layout
                if num_images == 1:
                    rows, cols = 1, 1
                elif num_images == 2:
                    rows, cols = 2, 1
                elif num_images <= 4:
                    rows, cols = 2, 2
                else:
                    rows, cols = 3, 2
                
                # cell dimensions with margins
                cell_w = area_w / cols
                cell_h = area_h / rows
                margin = Inches(0.1)
                
                for idx, image_path in enumerate(image_paths):
                    # Download if URL
                    if image_path.startswith('http'):
                        from utils.image_helper import download_image
                        cached_path = download_image(image_path)
                        if cached_path:
                            image_path = cached_path
                    
                    if image_path and os.path.exists(image_path):
                        try:
                            # Calculate grid position
                            row = idx // cols
                            col = idx % cols
                            
                            # Center point of the cell
                            center_x = area_x + (col * cell_w) + (cell_w / 2)
                            center_y = area_y + (row * cell_h) + (cell_h / 2)
                            
                            # Max dimensions for this image
                            max_w = cell_w - margin
                            max_h = cell_h - margin
                            
                            # Get actual image size to preserve aspect ratio
                            with PILImage.open(image_path) as img:
                                img_w, img_h = img.size
                                aspect_ratio = img_w / img_h
                            
                            # Calculate dimensions fitting within max_w/max_h
                            # Try fitting to width
                            final_w = max_w
                            final_h = final_w / aspect_ratio
                            
                            # If too tall, fit to height instead
                            if final_h > max_h:
                                final_h = max_h
                                final_w = final_h * aspect_ratio
                            
                            # Position top-left corner
                            pos_x = center_x - (final_w / 2)
                            pos_y = center_y - (final_h / 2)
                            
                            slide.shapes.add_picture(image_path, pos_x, pos_y, width=final_w, height=final_h)
                            
                        except Exception as e:
                            logger.error(f"  Failed to add image {idx + 1}: {e}")
            except ImportError:
                 logger.error("PIL not installed, defaulting to basic image layout")
                 # Fallback to simple single image
                 if image_paths:
                     try:
                         slide.shapes.add_picture(image_paths[0], Inches(0.5), Inches(1.8), width=Inches(4.5))
                     except: pass

        
        # Details box (right side) - adjusted position for taller header and multi-budget
        details_y = image_area_y  # Match image area Y position
        details_box = slide.shapes.add_textbox(Inches(5.2), details_y, Inches(4.3), Inches(4.5) if is_multibudget else Inches(5.0))
        details_frame = details_box.text_frame
        details_frame.word_wrap = True
        
        # Product Details heading
        p = details_frame.paragraphs[0]
        p.text = "Product Details"
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = RGBColor(26, 54, 93)  # Navy
        p.space_after = Pt(8)
        
        # Description label
        desc_label = details_frame.add_paragraph()
        desc_label.text = "Description:"
        desc_label.font.size = Pt(14)
        desc_label.font.bold = True
        desc_label.font.color.rgb = RGBColor(26, 54, 93)
        desc_label.space_after = Pt(4)
        
        # Full description paragraph - dynamically adjust font size based on length
        desc_p = details_frame.add_paragraph()
        # Truncate description if extremely long for PPT
        clean_desc = item['description']
        if len(clean_desc) > 800:
             clean_desc = clean_desc[:797] + "..."
        desc_p.text = clean_desc
        
        # Adjust font size based on description length
        desc_length = len(clean_desc)
        if desc_length > 600:
            desc_p.font.size = Pt(9)
        elif desc_length > 400:
            desc_p.font.size = Pt(10)
        elif desc_length > 200:
            desc_p.font.size = Pt(11)
        else:
            desc_p.font.size = Pt(12)
        
        desc_p.font.color.rgb = RGBColor(51, 51, 51)  # Dark text
        desc_p.space_after = Pt(8)
        
        # Key Details - professional text only (no icons)
        p = details_frame.add_paragraph()
        p.text = f"Brand: {item['brand']}"
        p.font.size = Pt(12)
        p.font.bold = True
        p.space_after = Pt(4)
        
        p = details_frame.add_paragraph()
        p.text = f"Quantity: {item['qty']} {item['unit']}"
        p.font.size = Pt(12)
        p.space_after = Pt(10)
        
        # Specifications heading
        p = details_frame.add_paragraph()
        p.text = "Specifications:"
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = RGBColor(26, 54, 93)
        p.space_after = Pt(6)
        
        # Add specifications (limit based on description length to fit page)
        max_specs = 3 if desc_length > 500 else 4 if desc_length > 300 else 5
        
        for spec in item['specifications'][:max_specs]:
            p = details_frame.add_paragraph()
            p.text = f"• {spec}"
            p.font.size = Pt(10 if desc_length > 500 else 11)
            p.font.color.rgb = RGBColor(51, 51, 51)
            p.level = 1
            p.space_after = Pt(2)
        
        # Warranty section (bottom left area) - SMALLER TEXT
        warranty_box = slide.shapes.add_textbox(Inches(0.6), Inches(6.5), Inches(4.0), Inches(0.5))
        warranty_frame = warranty_box.text_frame
        warranty_frame.word_wrap = True
        
        # Warranty heading
        warranty_title = warranty_frame.paragraphs[0]
        warranty_title.text = "Warranty"
        warranty_title.font.size = Pt(10) # Smaller
        warranty_title.font.bold = True
        warranty_title.font.color.rgb = RGBColor(26, 54, 93)
        warranty_title.space_after = Pt(2)
        
        # Warranty content
        warranty_content = warranty_frame.add_paragraph()
        warranty_content.text = "As per manufacturer - 5 years"
        warranty_content.font.size = Pt(9) # Smaller text
        warranty_content.font.color.rgb = RGBColor(80, 80, 80)
        
        # Footer with website
        footer_box = slide.shapes.add_textbox(Inches(3), Inches(7), Inches(4), Inches(0.3))
        footer_frame = footer_box.text_frame
        footer_frame.text = "https://alshayaenterprises.com"
        footer_p = footer_frame.paragraphs[0]
        footer_p.font.size = Pt(10)
        footer_p.font.color.rgb = RGBColor(100, 116, 139)
        footer_p.alignment = PP_ALIGN.CENTER
    
    def create_item_page_pdf(self, item, page_num):
        """Create PDF page for one item"""
        story = []
        
        # Item title
        item_title = f"Item {page_num}: {item['description'][:80]}"
        story.append(Paragraph(item_title, self.item_title_style))
        story.append(Spacer(1, 0.3*inch))
        
        # For multi-budget: Add reference image in upper left
        is_multibudget = item.get('is_multibudget', False)
        reference_image_path = item.get('reference_image_path')
        if is_multibudget and reference_image_path:
            try:
                # Download if URL
                if reference_image_path.startswith('http'):
                    from utils.image_helper import download_image
                    cached_path = download_image(reference_image_path)
                    if cached_path:
                        reference_image_path = cached_path
                
                if reference_image_path and os.path.exists(reference_image_path):
                    # Small reference image in upper left
                    ref_img = RLImage(reference_image_path, width=1.0*inch, height=0.75*inch)
                    # Position it using a table
                    ref_table = Table([[ref_img, Paragraph("Reference Image", ParagraphStyle('RefLabel', fontSize=8, textColor=colors.grey))]], 
                                     colWidths=[1.0*inch, 1.5*inch])
                    ref_table.setStyle(TableStyle([
                        ('VALIGN', (0, 0), (-1, -1), 'TOP'),
                        ('ALIGN', (0, 0), (0, 0), 'LEFT'),
                    ]))
                    story.append(ref_table)
                    story.append(Spacer(1, 0.1*inch))
            except Exception as e:
                logger.warning(f"Could not add reference image to PDF: {e}")
        
        # Create two-column layout
        left_content = []
        right_content = []
        
        # Left: Selected product image (bigger for multi-budget)
        image_path = item.get('image_path')
        if image_path:
            # If it's a URL, download it first
            if image_path.startswith('http'):
                from utils.image_helper import download_image
                cached_path = download_image(image_path)
                if cached_path:
                    image_path = cached_path
            
            if image_path and os.path.exists(image_path):
                try:
                    img = RLImage(image_path, width=2.5*inch, height=2.5*inch)
                    left_content.append(img)
                except Exception as e:
                    left_content.append(Paragraph('[Image Not Available]', self.styles['Normal']))
            else:
                left_content.append(Paragraph('[Image Not Available]', self.styles['Normal']))
        else:
            left_content.append(Paragraph('[Image Not Available]', self.styles['Normal']))
        
        # Right: Details
        details_html = f"""
            <para>
                <b><font size="14" color="#1a365d">Product Details</font></b><br/>
                <br/>
                <b>Brand:</b> {item['brand']}<br/>
                <b>Quantity:</b> {item['qty']} {item['unit']}<br/>
                <b>Unit Rate:</b> {item['unit_rate']}<br/>
                <b>Total Amount:</b> {item['total']}<br/>
                <br/>
                <b><font color="#1a365d">Specifications:</font></b><br/>
            </para>
        """
        right_content.append(Paragraph(details_html, self.spec_text_style))
        
        # Specifications
        for spec in item['specifications']:
            right_content.append(Paragraph(f"• {spec}", self.spec_text_style))
            right_content.append(Spacer(1, 0.05*inch))
        
        # Two-column table
        data = [[left_content, right_content]]
        t = Table(data, colWidths=[3*inch, 3.5*inch])
        t.setStyle(TableStyle([
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('LEFTPADDING', (0, 0), (-1, -1), 10),
            ('RIGHTPADDING', (0, 0), (-1, -1), 10),
        ]))
        
        story.append(t)
        
        return story
    
    def parse_items_from_extraction(self, extraction_result, session, file_id):
        """
        Parse individual items from extraction result
        Returns: list of item dictionaries
        """
        items = []
        
        for layout_result in extraction_result.get('layoutParsingResults', []):
            markdown_text = layout_result.get('markdown', {}).get('text', '')
            images = layout_result.get('markdown', {}).get('images', {})
            
            # Parse tables from markdown
            table_rows = self.extract_table_rows(markdown_text)
            
            for row in table_rows:
                item = {
                    'sn': row.get('sn', row.get('sl.no', '')),
                    'description': row.get('description', row.get('item', '')),
                    'qty': row.get('qty', row.get('quantity', '')),
                    'unit': row.get('unit', ''),
                    'unit_rate': row.get('unit rate', row.get('unit price', row.get('rate', ''))),
                    'total': row.get('total', row.get('amount', '')),
                    'image': self.find_item_image(row, images),
                    'brand': self.extract_brand(row.get('description', '')),
                    'specifications': self.extract_specifications(row.get('description', ''))
                }
                items.append(item)
        
        return items
    
    def extract_table_rows(self, markdown_text):
        """Extract table rows from markdown text"""
        lines = markdown_text.split('\n')
        rows = []
        headers = []
        
        for line in lines:
            if '|' in line:
                cells = [cell.strip() for cell in line.split('|') if cell.strip()]
                
                if not headers and cells:
                    # First row is headers
                    headers = [h.lower() for h in cells]
                elif headers and cells and not all(c in '-: ' for c in ''.join(cells)):
                    # Data row
                    if len(cells) == len(headers):
                        row = dict(zip(headers, cells))
                        rows.append(row)
        
        return rows
    
    def find_item_image(self, row, images):
        """Find image associated with this item"""
        # Try to find image reference in row
        for key, value in row.items():
            if 'image' in str(key).lower():
                # Check if value references an image
                value_str = str(value) if value else ''
                for img_path, img_url in images.items():
                    img_path_str = str(img_path) if img_path else ''
                    img_url_str = str(img_url) if img_url else ''
                    if value_str in img_path_str or img_path_str in value_str:
                        return img_url_str
                
        # Return first available image if no specific match
        if images:
            first_img = list(images.values())[0]
            return str(first_img) if first_img else None
        
        return None
    
    def extract_brand(self, description):
        """Extract brand name from description (simple heuristic)"""
        # Common brand patterns - this is simplified
        brands = ['Sedus', 'Narbutas', 'Sokoa', 'B&T', 'Herman Miller', 'Steelcase', 'Haworth', 'Knoll']
        
        for brand in brands:
            if brand.lower() in description.lower():
                return brand
        
        # Try to extract first capitalized word
        words = description.split()
        for word in words:
            if word and word[0].isupper() and len(word) > 2:
                return word
        
        return 'Premium Brand'
    
    def extract_specifications(self, description):
        """Extract specifications from description"""
        # Split description into bullet points
        specs = []
        
        # Look for dimensions
        dimension_pattern = r'\d+\s*[xX×]\s*\d+\s*[xX×]?\s*\d*\s*(mm|cm|m|inch|in|")'
        dimensions = re.findall(dimension_pattern, description)
        if dimensions:
            specs.append(f"Dimensions: {', '.join(dimensions)}")
        
        # Look for materials
        materials = ['wood', 'metal', 'steel', 'aluminum', 'fabric', 'leather', 'plastic', 'glass', 'laminate']
        found_materials = [mat for mat in materials if mat in description.lower()]
        if found_materials:
            specs.append(f"Materials: {', '.join(found_materials).title()}")
        
        # Look for colors
        colors_list = ['black', 'white', 'grey', 'gray', 'brown', 'blue', 'red', 'green', 'beige']
        found_colors = [col for col in colors_list if col in description.lower()]
        if found_colors:
            specs.append(f"Available Colors: {', '.join(found_colors).title()}")
        
        if not specs:
            # Use description as-is if no specific specs found
            specs.append(description[:200])
        
        return specs
    
    def create_cover_page(self):
        """Create presentation cover page"""
        story = []
        
        # Centered large logo (if available)
        logo = self._get_logo_path()
        if logo and os.path.exists(logo):
            try:
                img = RLImage(logo, width=3.5*inch, height=3.5*inch)
                img.hAlign = 'CENTER'
                story.append(img)
                story.append(Spacer(1, 0.2*inch))
            except Exception:
                pass

        # Title
        title = Paragraph("TECHNICAL PROPOSAL", self.title_style)
        story.append(Spacer(1, 2*inch))
        story.append(title)
        story.append(Spacer(1, 0.5*inch))
        
        # Subtitle
        subtitle_style = ParagraphStyle(
            'Subtitle',
            parent=self.styles['Normal'],
            fontSize=16,
            textColor=colors.HexColor('#1a365d'),
            alignment=TA_CENTER
        )
        subtitle = Paragraph("Furniture, Fixtures & Equipment", subtitle_style)
        story.append(subtitle)
        story.append(Spacer(1, 1*inch))
        
        # Company info
        company_info = f"""
            <para align="center">
                <b>Prepared By:</b><br/>
                <font size="14" color="#667eea"><b>Your Company Name</b></font><br/>
                <br/>
                Date: {datetime.now().strftime('%B %d, %Y')}<br/>
            </para>
        """
        story.append(Paragraph(company_info, self.styles['Normal']))
        
        return story
    
    def create_item_page(self, item, page_num):
        """Create one page for an item with eye-catching design"""
        story = []
        
        # Item number and title
        item_title = f"Item {page_num}: {item['description'][:60]}"
        story.append(Paragraph(item_title, self.item_title_style))
        story.append(Spacer(1, 0.2*inch))
        
        # Create two-column layout using table
        left_content = []
        right_content = []
        
        # Left column - Image
        if item['image']:
            try:
                # For now, placeholder - in production, download and embed image
                img_placeholder = Paragraph(
                    f'<para align="center"><b>[Product Image]</b><br/>{item["image"][:50]}...</para>',
                    self.styles['Normal']
                )
                left_content.append(img_placeholder)
            except:
                left_content.append(Paragraph('[Image Not Available]', self.styles['Normal']))
        else:
            left_content.append(Paragraph('[Image Not Available]', self.styles['Normal']))
        
        # Right column - Specifications
        specs_html = f"""
            <para>
                <b><font size="14" color="#667eea">Product Details</font></b><br/>
                <br/>
                <b>Brand:</b> {item['brand']}<br/>
                <b>Quantity:</b> {item['qty']} {item['unit']}<br/>
                <b>Unit Rate:</b> {item['unit_rate']}<br/>
                <b>Total Amount:</b> {item['total']}<br/>
                <br/>
                <b><font color="#667eea">Specifications:</font></b><br/>
            </para>
        """
        right_content.append(Paragraph(specs_html, self.spec_text_style))
        
        # Add specifications as bullet points
        for spec in item['specifications']:
            spec_bullet = f"• {spec}"
            right_content.append(Paragraph(spec_bullet, self.spec_text_style))
            right_content.append(Spacer(1, 0.1*inch))
        
        # Additional info
        additional_info = """
            <para>
                <br/>
                <b><font color="#667eea">Additional Information:</font></b><br/>
                • Country of Origin: Various<br/>
                • Warranty: As per manufacturer standard<br/>
                • Lead Time: 4-6 weeks<br/>
                • Finish: As specified or equivalent<br/>
            </para>
        """
        right_content.append(Paragraph(additional_info, self.spec_text_style))
        
        # Create two-column table
        data = [[left_content, right_content]]
        
        col_widths = [3*inch, 3.5*inch]
        t = Table(data, colWidths=col_widths)
        t.setStyle(TableStyle([
            ('VALIGN', (0, 0), (-1, -1), 'TOP'),
            ('LEFTPADDING', (0, 0), (-1, -1), 10),
            ('RIGHTPADDING', (0, 0), (-1, -1), 10),
        ]))
        
        story.append(t)
        story.append(Spacer(1, 0.3*inch))
        
        # Bottom section - Key features
        features_title = Paragraph('<b><font size="12" color="#1a365d">KEY FEATURES</font></b>', self.styles['Normal'])
        story.append(features_title)
        story.append(Spacer(1, 0.1*inch))
        
        features = [
            "✓ Premium quality construction",
            "✓ Modern ergonomic design",
            "✓ Environmentally friendly materials",
            "✓ Easy maintenance and durability"
        ]
        
        features_text = '<br/>'.join(features)
        story.append(Paragraph(features_text, self.spec_text_style))
        
        return story
